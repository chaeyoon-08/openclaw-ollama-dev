#!/usr/bin/env python3
"""
PPT 템플릿 생성 스크립트 - gcube 스타일
다크: 네이비 배경 + 퍼플 포인트
라이트: 흰 배경 + 퍼플 포인트
"""
# ref: https://python-pptx.readthedocs.io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

PURPLE      = RGBColor(0x7C, 0x3A, 0xED)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT   = RGBColor(0xB0, 0xB0, 0xC8)
DARK_BG     = RGBColor(0x0D, 0x0D, 0x1A)
DARK_CARD   = RGBColor(0x16, 0x16, 0x2E)
LIGHT_TEXT  = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_CARD  = RGBColor(0xF5, 0xF3, 0xFF)
LIGHT_SUB   = RGBColor(0x6B, 0x6B, 0x8A)
LIGHT_CONT  = RGBColor(0xF9, 0xF9, 0xFF)

W = Inches(13.33)
H = Inches(7.5)

def add_rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def add_text(slide, text, x, y, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def make_title_slide(prs, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    is_dark = theme == "dark"
    bg = DARK_BG if is_dark else WHITE
    tc = WHITE if is_dark else LIGHT_TEXT
    sc = GRAY_TEXT if is_dark else LIGHT_SUB
    cc = DARK_CARD if is_dark else LIGHT_CARD
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = bg
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), PURPLE)
    add_text(slide, "AI Assistant", Inches(0.6), Inches(1.8), Inches(4), Inches(0.4), 13, color=PURPLE)
    add_text(slide, "제목을 입력하세요", Inches(0.6), Inches(2.3), Inches(8), Inches(1.6), 44, bold=True, color=tc)
    add_rect(slide, Inches(0.6), Inches(4.1), Inches(1.5), Pt(2), PURPLE)
    add_text(slide, "부제목 / 발표자 이름", Inches(0.6), Inches(4.3), Inches(7), Inches(0.6), 18, color=sc)
    add_rect(slide, Inches(0), Inches(7.1), W, Inches(0.4), cc)
    add_text(slide, "Powered by Clari AI", Inches(0.4), Inches(7.12), Inches(5), Inches(0.3), 11, color=sc)

def make_content_slide(prs, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    is_dark = theme == "dark"
    bg = DARK_BG if is_dark else WHITE
    tc = WHITE if is_dark else LIGHT_TEXT
    cc = DARK_CARD if is_dark else LIGHT_CARD
    co = DARK_CARD if is_dark else LIGHT_CONT
    ct = GRAY_TEXT if is_dark else LIGHT_TEXT
    sc = GRAY_TEXT if is_dark else LIGHT_SUB
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = bg
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PURPLE)
    add_rect(slide, Inches(0), Inches(0.06), W, Inches(1.3), cc)
    add_rect(slide, Inches(0.35), Inches(0.25), Inches(0.07), Inches(0.9), PURPLE)
    add_text(slide, "슬라이드 제목", Inches(0.55), Inches(0.28), Inches(11.5), Inches(0.85), 30, bold=True, color=tc)
    add_rect(slide, Inches(0.35), Inches(1.55), Inches(12.6), Inches(5.55), co)
    add_text(slide, "• 내용을 입력하세요\n• 두 번째 항목\n• 세 번째 항목", Inches(0.6), Inches(1.75), Inches(12.1), Inches(5.1), 18, color=ct)
    add_rect(slide, Inches(0), Inches(7.1), W, Inches(0.4), cc)
    add_text(slide, "Powered by Clari AI", Inches(0.4), Inches(7.12), Inches(5), Inches(0.3), 11, color=sc)

def generate(out_dir, skills_dir):
    for theme in ["dark", "light"]:
        prs = Presentation()
        prs.slide_width = W
        prs.slide_height = H
        make_title_slide(prs, theme)
        make_content_slide(prs, theme)
        path = os.path.join(skills_dir, f"template_{theme}.pptx")
        prs.save(path)
        print(f"템플릿 생성 완료: {path}")

if __name__ == "__main__":
    skills_dir = os.path.join(
        os.environ.get("OPENCLAW_WORKSPACE", os.path.expanduser("~/.openclaw/workspace")),
        "skills", "office-document-specialist-suite"
    )
    out_dir = os.environ.get("OUTPUT_DIR", "/workspace/work")
    generate(out_dir, skills_dir)
