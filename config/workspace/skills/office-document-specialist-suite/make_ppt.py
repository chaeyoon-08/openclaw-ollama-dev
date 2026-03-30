#!/usr/bin/env python3
# ref: https://python-pptx.readthedocs.io
import argparse
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PURPLE = RGBColor(0x7C, 0x3A, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT = RGBColor(0xB0, 0xB0, 0xC8)
LIGHT_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_SUB = RGBColor(0x6B, 0x6B, 0x8A)

def add_rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def add_text(slide, text, x, y, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def make_title_slide(prs, title, subtitle, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    is_dark = theme == "dark"
    bg = RGBColor(0x0D, 0x0D, 0x1A) if is_dark else RGBColor(0xFF, 0xFF, 0xFF)
    tc = WHITE if is_dark else LIGHT_TEXT
    sc = GRAY_TEXT if is_dark else LIGHT_SUB
    cc = RGBColor(0x16, 0x16, 0x2E) if is_dark else RGBColor(0xF5, 0xF3, 0xFF)
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = bg
    add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), PURPLE)
    add_text(slide, "AI Assistant", Inches(0.6), Inches(1.8), Inches(4), Inches(0.4), 13, color=PURPLE)
    add_text(slide, title, Inches(0.6), Inches(2.3), Inches(8), Inches(1.6), 44, bold=True, color=tc)
    add_rect(slide, Inches(0.6), Inches(4.1), Inches(1.5), Pt(2), PURPLE)
    add_text(slide, subtitle, Inches(0.6), Inches(4.3), Inches(7), Inches(0.6), 18, color=sc)
    add_rect(slide, Inches(0), Inches(7.1), Inches(13.33), Inches(0.4), cc)
    add_text(slide, "Powered by Clari AI", Inches(0.4), Inches(7.12), Inches(5), Inches(0.3), 11, color=sc)

def make_content_slide(prs, title, content, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    is_dark = theme == "dark"
    bg = RGBColor(0x0D, 0x0D, 0x1A) if is_dark else RGBColor(0xFF, 0xFF, 0xFF)
    tc = WHITE if is_dark else LIGHT_TEXT
    cc = RGBColor(0x16, 0x16, 0x2E) if is_dark else RGBColor(0xF5, 0xF3, 0xFF)
    co = RGBColor(0x16, 0x16, 0x2E) if is_dark else RGBColor(0xF9, 0xF9, 0xFF)
    ct = GRAY_TEXT if is_dark else LIGHT_TEXT
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = bg
    add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), PURPLE)
    add_rect(slide, Inches(0), Inches(0.06), Inches(13.33), Inches(1.3), cc)
    add_rect(slide, Inches(0.35), Inches(0.25), Inches(0.07), Inches(0.9), PURPLE)
    add_text(slide, title, Inches(0.55), Inches(0.28), Inches(11.5), Inches(0.85), 30, bold=True, color=tc)
    add_rect(slide, Inches(0.35), Inches(1.55), Inches(12.6), Inches(5.55), co)
    add_text(slide, content, Inches(0.6), Inches(1.75), Inches(12.1), Inches(5.1), 18, color=ct)
    add_rect(slide, Inches(0), Inches(7.1), Inches(13.33), Inches(0.4), cc)
    add_text(slide, "Powered by Clari AI", Inches(0.4), Inches(7.12), Inches(5), Inches(0.3), 11,
             color=GRAY_TEXT if is_dark else LIGHT_SUB)

def make_ppt(slides_json, output, theme="dark"):
    os.makedirs(os.path.dirname(os.path.abspath(output)), exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    data = json.loads(slides_json)
    make_title_slide(prs, data.get("title", "제목"), data.get("subtitle", "Clari AI"), theme)
    for s in data.get("slides", []):
        make_content_slide(prs, s["title"], s["content"], theme)
    prs.save(output)
    print(f"PPT 생성 완료: {output}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--json", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument("--theme", default="dark", choices=["dark", "light"])
    args = parser.parse_args()
    make_ppt(args.json, args.output, args.theme)
